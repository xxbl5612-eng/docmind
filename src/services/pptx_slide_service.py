"""Extract slide structure from PPTX files for frontend rendering."""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from typing import Any


@dataclass
class SlideShapeData:
    shape_idx: int
    shape_type: str  # "text", "picture", "table", "group", "other"
    left: float
    top: float
    width: float
    height: float
    text: str | None = None
    font_size: float | None = None
    font_name: str | None = None
    font_bold: bool = False
    font_italic: bool = False
    font_color: str | None = None  # #RRGGBB
    fill_color: str | None = None
    alignment: str | None = None
    has_image: bool = False
    image_index: int | None = None
    table_rows: list[list[str]] | None = None
    paragraphs: list[dict[str, Any]] = field(default_factory=list)
    is_title: bool = False
    # Shape styling
    fill_type: str | None = None
    gradient_angle: float | None = None
    gradient_stops: list[dict] = field(default_factory=list)
    border_color: str | None = None
    border_width: float | None = None
    border_style: str | None = None
    border_radius: float | None = None
    shadow: bool = False
    rotation: float | None = None
    # Enhanced table
    table_data: dict | None = None


@dataclass
class SlideData:
    slide_index: int
    width_emu: int
    height_emu: int
    width_px: float
    height_px: float
    shapes: list[SlideShapeData] = field(default_factory=list)
    bg_color: str | None = None
    bg_fill_type: str | None = None
    bg_gradient_stops: list[dict] = field(default_factory=list)
    bg_gradient_angle: float | None = None


@dataclass
class PptxSlideExtract:
    slides: list[SlideData]
    image_count: int


EMU_TO_PX = 96.0 / 914400.0
FALLBACK_WIDTH_EMU = 12192000
FALLBACK_HEIGHT_EMU = 6858000
MAX_SLIDE_WIDTH_PX = 960

# Default theme colors (approximate Office theme)
THEME_COLORS: dict[str, str] = {
    "tx1": "000000", "tx2": "1F497D", "bg1": "FFFFFF", "bg2": "EEECE1",
    "accent1": "4F81BD", "accent2": "C0504D", "accent3": "9BBB59",
    "accent4": "8064A2", "accent5": "4BACC6", "accent6": "F79646",
}


def _resolve_color(color_obj) -> str | None:
    """Resolve a pptx color to hex string, including theme colors."""
    try:
        if color_obj is None:
            return None
        if hasattr(color_obj, "rgb") and color_obj.rgb:
            return str(color_obj.rgb)
        if hasattr(color_obj, "theme_color") and color_obj.theme_color:
            tc = str(color_obj.theme_color).lower()
            return THEME_COLORS.get(tc, "333333")
    except Exception:
        pass
    return None


def _get_effective_font(run, shape_is_title: bool) -> dict[str, Any]:
    """Get font properties from a run, with fallback defaults."""
    font = run.font
    result: dict[str, Any] = {}

    # Font size with fallback based on role
    if font.size:
        result["font_size"] = round(font.size / 12700, 1)
    else:
        result["font_size"] = 28 if shape_is_title else 18

    result["bold"] = font.bold if font.bold is not None else shape_is_title
    result["italic"] = bool(font.italic)

    name = font.name
    result["font_name"] = name if name else None

    color = _resolve_color(font.color)
    result["color"] = f"#{color}" if color else None

    return result


def _get_shape_fill_detail(shape) -> dict:
    """Extract fill type, solid color, and gradient stops."""
    result = {"fill_type": None, "fill_color": None, "gradient_angle": None, "gradient_stops": []}
    try:
        fill = shape.fill
        ft_val = fill.type if hasattr(fill, "type") else None
        ft = ft_val if isinstance(ft_val, str) else str(ft_val) if ft_val is not None else ""
        if "SOLID" in ft or (not ft and fill.fore_color):
            result["fill_type"] = "solid"
            fc = _resolve_color(fill.fore_color)
            if fc:
                result["fill_color"] = f"#{fc}"
        elif ft and ("GRADIENT" in ft.upper() or "2" in ft):
            result["fill_type"] = "gradient"
            try:
                for gs in fill.gradient_stops:
                    color = _resolve_color(gs.color)
                    if color:
                        result["gradient_stops"].append({"color": f"#{color}", "position": round(gs.position, 2)})
                result["gradient_angle"] = round(fill.gradient_angle or 0, 1)
            except Exception:
                pass
    except Exception:
        pass
    return result


def _get_slide_bg_detail(slide) -> dict:
    """Extract slide background including gradient."""
    result = {"bg_color": None, "bg_fill_type": None, "bg_gradient_stops": [], "bg_gradient_angle": None}
    try:
        bg = slide.background
        if bg and bg.fill:
            ft_val = bg.fill.type if hasattr(bg.fill, "type") else None
            ft = ft_val if isinstance(ft_val, str) else str(ft_val) if ft_val is not None else ""
            if "SOLID" in ft or (not ft and bg.fill.fore_color):
                result["bg_fill_type"] = "solid"
                fc = _resolve_color(bg.fill.fore_color)
                if fc:
                    result["bg_color"] = f"#{fc}"
            elif ft and ("GRADIENT" in ft.upper() or "2" in ft):
                result["bg_fill_type"] = "gradient"
                try:
                    for gs in bg.fill.gradient_stops:
                        color = _resolve_color(gs.color)
                        if color:
                            result["bg_gradient_stops"].append({"color": f"#{color}", "position": round(gs.position, 2)})
                    result["bg_gradient_angle"] = round(bg.fill.gradient_angle or 0, 1)
                except Exception:
                    pass
    except Exception:
        pass
    return result


def _get_shape_border(shape) -> dict:
    """Extract shape border/line properties."""
    result = {"border_color": None, "border_width": None, "border_style": None, "border_radius": None}
    try:
        line = shape.line
        if line and line.fill:
            fc = _resolve_color(line.fill.fore_color)
            if fc:
                result["border_color"] = f"#{fc}"
        if line and line.width:
            result["border_width"] = round(line.width / 12700, 1)
        if line and line.dash_style is not None:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            dash_map = {
                MSO_LINE_DASH_STYLE.SOLID: "solid",
                MSO_LINE_DASH_STYLE.DASH: "dashed",
                MSO_LINE_DASH_STYLE.DOT: "dotted",
                MSO_LINE_DASH_STYLE.DASH_DOT: "dashed",
            }
            result["border_style"] = dash_map.get(line.dash_style, "solid")
        # Rounded corners via XML
        from lxml import etree
        sp = shape._element
        prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is not None:
            prst = prstGeom.get('prst', '')
            if prst in ('roundRect', 'snipRoundRect', 'roundSameRect'):
                avLst = prstGeom.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if avLst is not None:
                    gd = avLst.find('{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
                    if gd is not None and gd.get('fmla'):
                        result["border_radius"] = round(int(gd.get('fmla').replace('val ', '')) / 12700, 1)
                    else:
                        result["border_radius"] = 8
                else:
                    result["border_radius"] = 8
    except Exception:
        pass
    return result


def _get_shape_effects(shape) -> dict:
    """Extract shadow and rotation from shape."""
    result = {"shadow": False, "rotation": None}
    try:
        if shape.rotation:
            result["rotation"] = round(shape.rotation, 1)
        from lxml import etree
        sp = shape._element
        effectLst = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
        if effectLst is not None:
            outerShdw = effectLst.find('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw')
            if outerShdw is not None:
                result["shadow"] = True
    except Exception:
        pass
    return result


def _extract_table_data(shape) -> dict:
    """Extract enhanced table data with merged cells and styles."""
    table = shape.table
    rows: list[list[str]] = []
    cell_styles: list[dict] = []
    row_count = len(table.rows)
    col_count = len(table.columns)

    total_w = sum(c.width for c in table.columns) or 1
    col_widths = [round(c.width / total_w, 3) for c in table.columns]

    header_count = 0
    if row_count > 0:
        first_row = table.rows[0]
        if any(cell.text and cell.text.strip() for cell in first_row.cells):
            header_count = 1

    for ri, row in enumerate(table.rows):
        row_texts: list[str] = []
        for ci, cell in enumerate(row.cells):
            row_texts.append(cell.text)
            style: dict = {"row": ri, "col": ci, "bg_color": None, "bold": False, "align": "left", "colspan": 1, "rowspan": 1}
            try:
                if cell.fill and cell.fill.fore_color:
                    fc = _resolve_color(cell.fill.fore_color)
                    if fc:
                        style["bg_color"] = f"#{fc}"
            except Exception:
                pass
            try:
                if cell.text_frame and cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        style["bold"] = bool(para.runs[0].font.bold)
                    style["align"] = _get_alignment(para)
            except Exception:
                pass
            try:
                if cell.is_merge_origin:
                    style["rowspan"] = cell.span_height or 1
                    style["colspan"] = cell.span_width or 1
                if cell.is_spanned:
                    style["colspan"] = 0
                    style["rowspan"] = 0
            except Exception:
                pass
            cell_styles.append(style)
        rows.append(row_texts)

    return {
        "rows": rows, "col_widths": col_widths,
        "header_count": header_count, "cell_styles": cell_styles,
        "row_count": row_count, "col_count": col_count,
    }


def _get_alignment(para) -> str:
    try:
        from pptx.enum.text import PP_ALIGN
        mapping = {
            PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify",
        }
        return mapping.get(para.alignment, "left") if para.alignment else "left"
    except Exception:
        return "left"


def _extract_image(shape) -> tuple[bytes, str] | None:
    """Extract image blob and content_type from a shape (picture, fill, etc.)."""
    # Direct picture
    if hasattr(shape, "image") and shape.image:
        img = shape.image
        return img.blob, img.content_type or "image/png"

    # Fill image (e.g. shape background)
    try:
        fill = shape.fill
        if fill and hasattr(fill, "type") and str(fill.type) == "PATTERN":
            pass  # skip patterns
    except Exception:
        pass

    return None


def _is_group(shape) -> bool:
    """Check if a shape is a group shape (has sub-shapes)."""
    try:
        st = shape.shape_type
        # MSO_SHAPE_TYPE.GROUP = 6
        if st is not None and int(st) == 6:
            return True
    except Exception:
        pass
    # Fallback: check if it has a 'shapes' attribute with sub-elements
    if hasattr(shape, "shapes"):
        try:
            sub = shape.shapes
            if sub is not None and len(sub) > 0:
                return True
        except Exception:
            pass
    return False


def _apply_shape_styling(shape, sh_data: SlideShapeData) -> None:
    """Apply border, shadow, rotation, and fill styling to a shape data object."""
    border = _get_shape_border(shape)
    sh_data.border_color = border["border_color"]
    sh_data.border_width = border["border_width"]
    sh_data.border_style = border["border_style"]
    sh_data.border_radius = border["border_radius"]
    effects = _get_shape_effects(shape)
    sh_data.shadow = effects["shadow"]
    sh_data.rotation = effects["rotation"]
    # Fill for non-text shapes (text shapes handle fill in _process_text_shape)
    if sh_data.shape_type in ("picture", "table", "other"):
        fill = _get_shape_fill_detail(shape)
        sh_data.fill_type = fill["fill_type"]
        if fill["fill_type"] == "solid":
            sh_data.fill_color = fill["fill_color"]
        elif fill["fill_type"] == "gradient":
            sh_data.gradient_stops = fill["gradient_stops"]
            sh_data.gradient_angle = fill["gradient_angle"]


def _count_images_recursive(shapes, slide_img_idx: int, result_shapes: list, scale: float) -> tuple[int, list]:
    """Recursively extract shapes including from groups. Returns (next_image_idx, shapes)."""

    next_idx = slide_img_idx
    all_shapes: list = []

    for sh_idx, shape in enumerate(shapes):
        # Check if this is a group shape
        if _is_group(shape):
            # Recursively process group shapes
            sub_idx, sub_shapes = _count_images_recursive(shape.shapes, next_idx, [], scale)
            next_idx = sub_idx
            # Adjust positions to be relative to the group's position
            for sub in sub_shapes:
                sub.left += round(shape.left * EMU_TO_PX * scale, 1)
                sub.top += round(shape.top * EMU_TO_PX * scale, 1)
            all_shapes.extend(sub_shapes)
            continue

        sh_data = SlideShapeData(
            shape_idx=sh_idx,
            shape_type="other",
            left=round(shape.left * EMU_TO_PX * scale, 1),
            top=round(shape.top * EMU_TO_PX * scale, 1),
            width=round(shape.width * EMU_TO_PX * scale, 1),
            height=round(shape.height * EMU_TO_PX * scale, 1),
        )

        img = _extract_image(shape)
        if img is not None:
            sh_data.shape_type = "picture"
            sh_data.has_image = True
            sh_data.image_index = next_idx
            next_idx += 1
            _apply_shape_styling(shape, sh_data)
            all_shapes.append(sh_data)
            continue

        if shape.has_table:
            sh_data.shape_type = "table"
            sh_data.table_data = _extract_table_data(shape)
            _apply_shape_styling(shape, sh_data)
            all_shapes.append(sh_data)
            continue

        _process_text_shape(shape, sh_data)
        _apply_shape_styling(shape, sh_data)
        all_shapes.append(sh_data)

    return next_idx, all_shapes


def _process_text_shape(shape, sh_data: SlideShapeData) -> None:
    """Fill in text-related fields on sh_data from a shape."""
    if not shape.has_text_frame:
        return

    sh_data.shape_type = "text"
    role = _shape_role(shape)
    sh_data.is_title = (role == "title")
    tf = shape.text_frame
    paragraphs: list[dict[str, Any]] = []
    all_text: list[str] = []

    for para in tf.paragraphs:
        para_runs: list[dict[str, Any]] = []
        para_text_parts: list[str] = []

        for run in para.runs:
            font_info = _get_effective_font(run, role == "title")
            run_data: dict[str, Any] = {"text": run.text, **font_info}
            para_runs.append(run_data)
            para_text_parts.append(run.text)

        para_text = "".join(para_text_parts)
        if para_text.strip():
            all_text.append(para_text)

        # Bullet detection
        bullet_char = None
        bullet_type = "none"
        try:
            pPr = para._element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            if pPr is not None:
                buChar = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                if buChar is not None:
                    bullet_char = buChar.get('char')
                    bullet_type = "bullet"
                else:
                    buNone = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                    if buNone is None and para.level and para.level > 0:
                        bullet_type = "bullet"
                        bullet_char = "●" if para.level == 1 else "○"
        except Exception:
            pass

        paragraphs.append({
            "text": para_text,
            "runs": para_runs,
            "alignment": _get_alignment(para),
            "level": para.level if para.level else 0,
            "bullet_type": bullet_type,
            "bullet_char": bullet_char,
        })

    sh_data.paragraphs = paragraphs
    sh_data.text = "\n".join(all_text)

    for para in paragraphs:
        if para["runs"]:
            first = para["runs"][0]
            sh_data.font_size = first.get("font_size")
            sh_data.font_name = first.get("font_name")
            sh_data.font_bold = first.get("bold", False)
            sh_data.font_italic = first.get("italic", False)
            sh_data.font_color = first.get("color")
            sh_data.alignment = para["alignment"]
            break

    if sh_data.font_size is None:
        sh_data.font_size = 28 if role == "title" else 18
    if role == "title" and not sh_data.font_bold:
        sh_data.font_bold = True

    fill = _get_shape_fill_detail(shape)
    sh_data.fill_type = fill["fill_type"]
    if fill["fill_type"] == "solid":
        sh_data.fill_color = fill["fill_color"]
    elif fill["fill_type"] == "gradient":
        sh_data.gradient_stops = fill["gradient_stops"]
        sh_data.gradient_angle = fill["gradient_angle"]


def _shape_role(shape) -> str:
    """Determine if a shape is a title, subtitle, or body."""
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type is not None:
                from pptx.enum.shapes import PP_PLACEHOLDER
                if ph.type == PP_PLACEHOLDER.TITLE:
                    return "title"
                if ph.type == PP_PLACEHOLDER.SUBTITLE:
                    return "subtitle"
                if ph.type == PP_PLACEHOLDER.CENTER_TITLE:
                    return "title"
        # Heuristic: check shape name
        name = shape.name.lower() if hasattr(shape, "name") else ""
        if "title" in name:
            return "title"
        if "subtitle" in name:
            return "subtitle"
    except Exception:
        pass
    return "body"


def extract_slides(data: bytes) -> PptxSlideExtract:
    """Extract slide structure from PPTX bytes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slide_width = prs.slide_width or FALLBACK_WIDTH_EMU
    slide_height = prs.slide_height or FALLBACK_HEIGHT_EMU
    scale = min(MAX_SLIDE_WIDTH_PX / (slide_width * EMU_TO_PX), 1.0)

    slides: list[SlideData] = []
    total_images = 0

    for s_idx, slide in enumerate(prs.slides):
        bg_detail = _get_slide_bg_detail(slide)
        sd = SlideData(
            slide_index=s_idx,
            width_emu=slide_width,
            height_emu=slide_height,
            width_px=round(slide_width * EMU_TO_PX * scale, 1),
            height_px=round(slide_height * EMU_TO_PX * scale, 1),
            bg_color=bg_detail["bg_color"],
            bg_fill_type=bg_detail["bg_fill_type"],
            bg_gradient_stops=bg_detail["bg_gradient_stops"],
            bg_gradient_angle=bg_detail["bg_gradient_angle"],
        )
        slide_image_idx, slide_shapes = _count_images_recursive(
            slide.shapes, 0, [], scale
        )
        total_images += slide_image_idx
        sd.shapes = slide_shapes
        slides.append(sd)

    return PptxSlideExtract(slides=slides, image_count=total_images)


def extract_slide_image(data: bytes, slide_index: int, image_index: int) -> tuple[bytes, str] | None:
    """Extract a specific embedded image from a PPTX file (per-slide index)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    if slide_index >= len(prs.slides):
        return None

    slide = prs.slides[slide_index]

    def _find_in_shapes(shapes, target_idx: int, counter: list[int]):
        for shape in shapes:
            # Check group shapes recursively
            if _is_group(shape):
                result = _find_in_shapes(shape.shapes, target_idx, counter)
                if result is not None:
                    return result
                continue

            img = _extract_image(shape)
            if img is not None:
                if counter[0] == target_idx:
                    return img
                counter[0] += 1
        return None

    result = _find_in_shapes(slide.shapes, image_index, [0])
    return result
