"""Document parsing service: extract text from various file formats."""

from __future__ import annotations

import io


def parse_pdf(data: bytes) -> str:
    """Extract text from PDF using pdfplumber, falling back to OCR for image-based PDFs."""
    import pdfplumber
    parts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    result = "\n\n".join(parts).strip()

    # If no text found, try OCR on rendered pages (image-based PDF)
    if not result:
        try:
            from src.ai.ocr import ocr_pdf_pages
            result = ocr_pdf_pages(data)
        except Exception:
            pass

    return result


def parse_docx(data: bytes) -> str:
    """Extract text from DOCX."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def parse_pptx(data: bytes) -> str:
    """Extract text from PPTX."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
        if slide_texts:
            parts.append("\n".join(slide_texts))
    return "\n\n---\n\n".join(parts)


def parse_xlsx(data: bytes) -> str:
    """Extract data from XLSX as TSV-like text."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def parse_csv(data: bytes) -> str:
    """Parse CSV and return as text representation."""
    text = data.decode("utf-8", errors="replace")
    return text


def parse_text(data: bytes) -> str:
    """Parse plain text with encoding detection."""
    from src.utils.text_utils import decode_bytes
    return decode_bytes(data)


def parse_markdown(data: bytes) -> str:
    """Parse Markdown (identity transform)."""
    return data.decode("utf-8", errors="replace")


def parse_html(data: bytes) -> str:
    """Extract text from HTML."""
    from src.ai.standardizer import _html_to_plain
    raw = data.decode("utf-8", errors="replace")
    return _html_to_plain(raw)


def parse_doc(data: bytes) -> str:
    """Extract text from legacy .doc files (limited support).

    .doc is a proprietary binary format. For best results, convert .doc
    files to .docx before uploading. Falls back to plain-text decoding
    which may produce partial extraction with binary noise.
    """
    return parse_text(data)


def parse_image(data: bytes) -> str:
    """Extract text from images (PNG/JPEG) using OCR."""
    from src.ai.ocr import ocr_image
    return ocr_image(data)


# ── Parser registry ──

PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "doc": parse_doc,
    "pptx": parse_pptx,
    "xlsx": parse_xlsx,
    "csv": parse_csv,
    "txt": parse_text,
    "md": parse_markdown,
    "html": parse_html,
    "png": parse_image,
    "jpg": parse_image,
    "jpeg": parse_image,
}


def parse_document(data: bytes, input_format: str) -> str:
    """Parse document bytes into plain text using the appropriate parser."""
    parser = PARSERS.get(input_format)
    if parser is None:
        # Fallback: try to decode as text
        from src.utils.text_utils import decode_bytes
        return decode_bytes(data)
    return parser(data)
